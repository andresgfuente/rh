import boto3, json
import os
from langchain.schema import Document
from pypdf import PdfReader
import boto3
from langchain_aws import ChatBedrock
from langchain.output_parsers import ResponseSchema, StructuredOutputParser
from langchain import PromptTemplate
import re
import ast
from docx import Document as DocxDocument 
import time
import streamlit as st
from io import BytesIO

from pptx import Presentation

session = boto3.Session()

s3_client = boto3.client('s3', region_name='us-east-1')
textract_client = boto3.client('textract',region_name='us-east-1')

# Configurar el nombre de tu bucket
bucket_name = 'textract-console-us-east-1-4a58348c-7a8e-4dc3-b4d9-9478d91fe43f'

session = boto3.Session()
bedrock = session.client(service_name='bedrock-runtime',region_name='us-east-1')
#Extract Information from PDF file
def read_pdf_data(pdf_file):
    pdf_page = PdfReader(pdf_file)
    text = ""
    for page in pdf_page.pages:
        text += page.extract_text()
    return text

def clasificador_texto(texto):
    session=boto3.Session()
    bedrock=session.client(service_name='bedrock-runtime',region_name='us-east-1')
    llm=ChatBedrock(client=bedrock,model_id='anthropic.claude-3-sonnet-20240229-v1:0',region_name='us-east-1',model_kwargs={"temperature": 0.0})

    response_schema=[
        ResponseSchema(name='Nombre',description='Nombre del Candidato. Por ejemplo: Raul, Juan, Oscar'),
        ResponseSchema(name='Apellido',description='Apellido del Candidato. Por ejemplo: Gomez, Salazar, Baez.'),
        ResponseSchema(name='Nombre y Apellido',description='Nombre y Apellido del Candidato. Por ejemplo: Raul Gomez'),
        ResponseSchema(name='Sexo',description='Sexo del Candidato. Las opciones son: Masculino, Femenino'),
        ResponseSchema(name='Correo',description='Correo electronico del Candidato. Responde No se menciona en caso de no encontrar la respuesta'),
        ResponseSchema(name='Telefono',description='Telefono del Candidato. Responde No se menciona en caso de no encontrar la respuesta'),
        ResponseSchema(name='Fecha de Nacimiento',description='Fecha de Nacimiento del Candidato. Responde No se menciona en caso de no encontrar la respuesta'),
        ResponseSchema(name='Formacion Academica',description='Cual es la formacion academica del candidato? Sea breve y conciso en la respuesta, citando los datos mas relevantes.'),
        ResponseSchema(name='Experiencia Laboral',description='Cual es la experiencia laboral?Sea breve y conciso en la respuesta, citando los datos mas relevantes.'),
        ResponseSchema(name='Fortalezas',description='Cuales son las Fortalezas del Candidato con respecto al la descripcion del trabajo?'),
        ResponseSchema(name='Debilidades',description='Cuales son las Debilidades del Candidato con respecto al la descripcion del trabajo?'),
        ResponseSchema(name='Categoria',description='En que Area de la organizacion puede el candidato encajar basado en sus aptitudes y habilidades? No agregues informacion extra. Expresate con maximo 5 palabras')
    ]

    output_parser=StructuredOutputParser.from_response_schemas(response_schema)
    format_instructions=output_parser.get_format_instructions()
    template=""" A continuacion, se te proporcionara texto sobre un curriculum vitae. En base al texto proporcionado deberas extraer la siguiente informacion con el siguiente formato:
    {format_instructions} 
    {texto}


    """

    prompt=PromptTemplate(template=template, input_variables=[texto],partial_variables={'format_instructions':format_instructions})
    #prompt=prompt.format_prompt(texto=texto).text
    #respuesta=llm.invoke(prompt)
    #st.write(respuesta)
    #pattern=r'({[^}]+})'
    #matches=re.findall(pattern,respuesta)
    #st.write(matches)
    chain=prompt|llm|output_parser
    respuesta=chain.invoke(texto)
    #respuesta=ast.literal_eval(matches[0])
    #chain=prompt|llm|output_parser
    #respuesta=chain.invoke(texto)
    return respuesta

# iterate over files in 
# that user uploaded PDF files, one by one
def create_docs(user_file_list, unique_id):
    docs = []
    
    for file in user_file_list:
        name=f'raw/cvs/{file.name}'
        if file is not None:
             # Restablecer puntero del archivo original si necesitas usarlo de nuevo
            
        # Leer contenido dependiendo del tipo de archivo
            if file.name.endswith(".pdf"):
                
                chunks = read_pdf_data(file)  # Función para extraer texto de PDFs
            elif file.name.endswith(".docx"):
                chunks = read_docx_data(file)  # Nueva función para extraer texto de Word
            elif file.name.lower().endswith((".png", ".jpg", ".jpeg")):
                name = f'raw/cvs/{file.name}'
                file.seek(0)  # Asegura que el puntero del archivo esté al inicio antes de subirlo
                upload_to_s3(file, bucket_name, name)
                file.seek(0)  # Restablece el puntero antes de procesarlo con Textract
                chunks = process_pdf(bucket=bucket_name, document=name) 
                st.write(chunks)# Nueva función para extraer texto de imágenes
            elif file.name.endswith(".pptx"):
                chunks = extract_text_from_pptx(file)
                #st.write(chunks)# Nueva función para extraer texto de PowerPoint
            else:
                raise ValueError(f"Tipo de archivo no soportado: {file.name}")

        if len(chunks.strip()) < 5:  # Si el texto es menor a 5 caracteres
            try:
                print(f"El contenido extraído de {file.name} es insuficiente. Intentando OCR...")
                print(f'realizando proceso de ocr para archivo {file.name.lower()}')
                #upload_to_s3(file, bucket_name, file.name)
                file.seek(0)  # Asegura que el puntero del archivo esté al inicio antes de subirlo
                upload_to_s3(file, bucket_name, name)
                file.seek(0)  # Restablece el puntero antes de procesarlo con Textract
                chunks = process_pdf(bucket=bucket_name, document=name)
                #st.write(chunks)
                if chunks==None or len(chunks)<=10:
                    print('tiene algun error')
                    #raise ValueError(f'Verificar el archivo {file.name.lower()}, cuenta con algun error')


            except Exception as e:
                st.error(f"Verificar manualmente, fallo el archivo: {file.name.lower()}: {e}")
                chunks = "[Contenido no extraído]"
                continue
        
        # Agregar elementos a la lista con los datos y metadatos
        
        docs.append(Document(
            page_content=chunks,
            metadata={
                "name": file.name,
                "id": file.file_id,
                "type": file.type,
                "size": file.size,
                "unique_id": unique_id,
            },
        ))

    return docs





def embedding(text):
    body = json.dumps({
            "inputText": text,
        })
    accept = "application/json"
    content_type = "application/json"

    response = bedrock.invoke_model(
        body=body, modelId="amazon.titan-embed-text-v1", accept=accept, contentType=content_type
    )

    response_body = json.loads(response.get('body').read())

    return response_body['embedding']


def fortalezas_y_debilidades(texto,descripcion_trabajo):
    session=boto3.Session()
    bedrock=session.client(service_name='bedrock-runtime',region_name='us-east-1')
    llm=ChatBedrock(client=bedrock,model_id='anthropic.claude-3-sonnet-20240229-v1:0',region_name='us-east-1',model_kwargs={"temperature": 0.1})

    response_schema=[
        ResponseSchema(name='Fortalezas',description='Cuales son las Fortalezas del Candidato con respecto al la descripcion del trabajo?'),
        ResponseSchema(name='Debilidades',description='Cuales son las Debilidades del Candidato con respecto al la descripcion del trabajo?'),
        ResponseSchema(name='Categoria',description='En que Area de la organizacion puede el candidato encajar basado en sus aptitudes y habilidades? No agregues informacion extra. Expresate con maximo 5 palabras'),
    ]

    output_parser=StructuredOutputParser.from_response_schemas(response_schema)
    format_instructions=output_parser.get_format_instructions()
    template=""" 
    Actua como un reclutador de talentos con experiencia en el ambito bancario. 
    En base a la descripcion del puesto de trabajo que se te proporcionara deberas realizar un analisis 
    con las siguientes instrucciones tomando en cuenta los datos del curriculum vitae. 
    Tu respuesta debe ser bien detallada y precisa.
   A continuacion se presentan las instrucciones:
    {format_instructions} \n
    La descripcion de la busqueda del trabajo son las siguientes:
   {descripcion_trabajo}
    A continuacion se presentan los datos del curriculum vitae de la persona:
    {texto}

    """

    prompt = PromptTemplate(
    template=template,
    input_variables=["descripcion_trabajo", "texto"],
    partial_variables={"format_instructions": format_instructions})
    #prompt=prompt.format_prompt(texto=texto).text
    #respuesta=llm.invoke(prompt)
    prompt_text = prompt.format(descripcion_trabajo=descripcion_trabajo, texto=texto)
    respuesta=llm.invoke(prompt_text).content
    pattern=r'({[^}]+})'
    matches=re.findall(pattern,respuesta)
    diccionario=ast.literal_eval(matches[0])
    fortalezas=diccionario['Fortalezas']
    debilidades=diccionario['Debilidades']
    categoria=diccionario['Categoria']
    return fortalezas,debilidades,categoria

def analizar_requerimientos_excluyentes(texto,requerimientos):
    session=boto3.Session()
    bedrock=session.client(service_name='bedrock-runtime',region_name='us-east-1')
    llm=ChatBedrock(client=bedrock,model_id='anthropic.claude-3-sonnet-20240229-v1:0',region_name='us-east-1',model_kwargs={"temperature": 0.1})
    template='''
    Actua como un reclutador de talentos con experiencia en el ambito bancario.
    A continuacion, se te proporcionara informacion sobre el curriculum vitae de un candidato 
    y tambien los requerimientos excluyentes de la busqueda de la persona que se desea contratar. 
    Deberas responder solamente con:
    Excluyente: en caso de que no cumpla con los requerimientos.
    No Excluyente: en caso de que cumpla con los requerimientos.
    Las opciones de respuesta son: Excluyente, No Excluyente
    Tambien debes agregar la razon de porque es No Excluyente y porque es Excluyente sea cual fuera el caso.
    No debes agregar texto extra a la respuesta, debes ser directo y conciso con la respuesta
    Los Requerimientos son los siguientes: \n 
    {requerimientos}
    La informacion del Curriculum Vitae del candidato es la siguiente: \n
    {texto}
    '''

    prompt = PromptTemplate(
    template=template,
    input_variables=["texto", "requerimientos"])
    prompt_text = prompt.format(requerimientos=requerimientos, texto=texto)
    respuesta=llm.invoke(prompt_text).content
    return respuesta



def upload_json_to_s3(data, bucket_name, file_name):
    """Sube un archivo JSON al bucket de S3."""
    json_buffer = BytesIO(json.dumps(data).encode('utf-8'))  # Convertir el JSON a bytes
    s3_client.put_object(Bucket=bucket_name, Key=file_name, Body=json_buffer.getvalue())
    return f"s3://{bucket_name}/{file_name}"



# Función para subir el archivo a S3
def upload_to_s3(file, bucket, object_name):
    try:
        file_bytes = BytesIO(file.read())  # Lee el contenido del archivo en memoria
        file_bytes.seek(0)  # Reinicia el puntero antes de subirlo
        s3_client.upload_fileobj(file_bytes, bucket, object_name)
        file.seek(0)  # Vuelve a reiniciar el puntero después de subirlo
    except Exception as e:
        st.error(f'Error al subir el archivo: {e}')


def start_textract(bucket, document):
    response = textract_client.start_document_analysis(
        DocumentLocation={
            'S3Object': {
                'Bucket': bucket,
                'Name': document,
            }
        },
        FeatureTypes=["TABLES", "FORMS"]  # Ajusta según lo que necesites
    )
    return response['JobId']

# Función para verificar el estado del trabajo
def check_job_status(job_id):
    response = textract_client.get_document_analysis(JobId=job_id)
    return response

# Función para procesar el PDF
def process_pdf(bucket, document):
    job_id = start_textract(bucket, document)
    print(f'Se inició el trabajo con ID: {job_id}')

    # Esperar hasta que el trabajo esté completo
    while True:
        response = check_job_status(job_id)
        status = response['JobStatus']
        print(f'Status del trabajo: {status}')

        if status in ['SUCCEEDED', 'FAILED']:
            break

        time.sleep(5)  # Espera antes de volver a verificar el estado

    if status == 'SUCCEEDED':
        # Recuperar y procesar todas las páginas usando NextToken
        text = ''
        next_token = None
        
        while True:
            if next_token:
                # Si hay más páginas, hacer otra llamada con el NextToken
                response = textract_client.get_document_analysis(
                    JobId=job_id,
                    NextToken=next_token
                )
            else:
                # Primera llamada ya realizada previamente en check_job_status
                response = check_job_status(job_id)

            # Procesar los resultados
            pages = response.get('Blocks', [])
            for block in pages:
                if block['BlockType'] == 'LINE':
                    text += block['Text'] + '\n'
            
            # Verificar si hay más resultados (páginas)
            next_token = response.get('NextToken')
            if not next_token:
                break
        
        
        return text
    
    
def read_docx_data(file):
    docx = DocxDocument(file)
    text = []

    # Extraer texto de párrafos
    for paragraph in docx.paragraphs:
        text.append(paragraph.text)

    # Extraer texto de tablas
    for table in docx.tables:
        for row in table.rows:
            for cell in row.cells:
                text.append(cell.text)

    return "\n".join(text)



def extract_text_from_pptx(file):
    """
    Extrae el texto de un archivo PowerPoint (.pptx).

    Args:
        file: Archivo PowerPoint a procesar.

    Returns:
        str: Texto extraído del PowerPoint.
    """
    try:
        # Cargar la presentación
        presentation = Presentation(file)
        text = []

        # Iterar a través de las diapositivas y sus elementos
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:  # Verificar si el elemento tiene un marco de texto
                    for paragraph in shape.text_frame.paragraphs:
                        text.append(paragraph.text)

        return "\n".join(text)
    except Exception as e:
        raise ValueError(f"Error al procesar el archivo PowerPoint: {e}")